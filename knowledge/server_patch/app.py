from __future__ import annotations

import asyncio
import base64
import hashlib
import json
import mimetypes
import os
import platform
import re
import requests
import shutil
import sqlite3
import subprocess
import tempfile
import time
import uuid
from contextlib import contextmanager
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Optional

import httpx
from fastapi import Depends, FastAPI, File, Form, HTTPException, Request, Security, UploadFile
from fastapi.responses import FileResponse, Response
from fastapi.security import HTTPAuthorizationCredentials, HTTPBearer
from pptx import Presentation
from pypdf import PdfReader
try:
    import fitz
except Exception:  # pragma: no cover - optional dependency fallback
    fitz = None

app = FastAPI(title="Knowledge File API")
security = HTTPBearer()

BASE_DIR = Path(
    os.getenv("OPENCLAW_KNOWLEDGE_BASE_DIR", str(Path(__file__).resolve().parent))
).resolve()
RUNTIME_ROOT = Path(
    os.getenv("OPENCLAW_RUNTIME_ROOT", str(BASE_DIR.parent.parent))
).resolve()
DATA_DIR = BASE_DIR / "data"
FILES_DIR = DATA_DIR / "files"
DERIVED_DIR = DATA_DIR / "derived"
DB_PATH = DATA_DIR / "kb.sqlite"

ZHIPU_API_KEY = os.getenv("ZHIPU_API_KEY", "").strip()
OPENAI_API_KEY = (os.getenv("OPENAI_API_KEY") or os.getenv("MEM0_OPENAI_API_KEY") or "").strip()
OPENAI_BASE_URL = os.getenv("OPENAI_BASE_URL", "https://api.openai.com/v1").rstrip("/")
OPENAI_VISION_MODEL = os.getenv("OPENAI_VISION_MODEL", "gpt-4.1-mini").strip() or "gpt-4.1-mini"
OPENAI_MAX_RETRIES = int(os.getenv("OPENAI_MAX_RETRIES", "4"))
OPENAI_RETRY_BACKOFF_MS = int(os.getenv("OPENAI_RETRY_BACKOFF_MS", "600"))
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "").strip()
GEMINI_BASE_URL = os.getenv(
    "GEMINI_BASE_URL",
    "https://generativelanguage.googleapis.com/v1beta/openai",
).rstrip("/")

EMBEDDING_PROVIDER = os.getenv(
    "EMBEDDING_PROVIDER",
    "zhipu" if ZHIPU_API_KEY else "openai",
).strip().lower()
DEFAULT_EMBEDDING_MODEL = "embedding-3" if EMBEDDING_PROVIDER == "zhipu" else "text-embedding-3-small"
DEFAULT_EMBEDDING_DIM = 2048 if EMBEDDING_PROVIDER == "zhipu" else 1536
EMBEDDING_MODEL = os.getenv("EMBEDDING_MODEL", DEFAULT_EMBEDDING_MODEL)
EMBEDDING_DIM = int(os.getenv("EMBEDDING_DIM", str(DEFAULT_EMBEDDING_DIM)))

QDRANT_URL = os.getenv("QDRANT_URL", "http://127.0.0.1:6333")
DEFAULT_QDRANT_COLLECTION = (
    "company-kb-files-v2" if EMBEDDING_PROVIDER == "zhipu" else "company-kb-local-v1"
)
QDRANT_COLLECTION = os.getenv("QDRANT_COLLECTION", DEFAULT_QDRANT_COLLECTION)
if QDRANT_COLLECTION == "company-kb-prod" and EMBEDDING_PROVIDER == "zhipu":
    QDRANT_COLLECTION = "company-kb-files-v2"
QDRANT_API_KEY = os.getenv("QDRANT_API_KEY")
API_TOKEN = os.getenv("API_TOKEN", "")

VISION_PROVIDER = os.getenv(
    "VISION_PROVIDER",
    "zhipu" if ZHIPU_API_KEY else "gemini",
).strip().lower()
VISION_MODEL = os.getenv(
    "VISION_MODEL",
    "glm-4.6v-flashx" if VISION_PROVIDER == "zhipu" else "gemini-2.5-pro",
)

IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}
PPT_SUFFIXES = {".pptx"}
LEGACY_PPT_SUFFIXES = {".ppt"}
PDF_SUFFIXES = {".pdf"}
VISION_OCR_SCRIPT = (
    RUNTIME_ROOT
    / "workspace"
    / "xiaoguan"
    / "wechat-ops-runtime"
    / "scripts"
    / "vision_ocr.swift"
)

def require_env(name: str, value: str) -> str:
    if value:
        return value
    raise HTTPException(status_code=500, detail=f"Missing required environment variable: {name}")


def qdrant_headers() -> dict[str, str]:
    headers = {"Content-Type": "application/json"}
    if QDRANT_API_KEY:
        headers["api-key"] = QDRANT_API_KEY
    return headers


def qdrant_api_url(path: str) -> str:
    return f"{QDRANT_URL.rstrip('/')}{path}"


def qdrant_request(
    method: str, path: str, payload: Optional[dict[str, Any]] = None
) -> dict[str, Any]:
    try:
        with httpx.Client(timeout=30.0) as client:
            response = client.request(
                method,
                qdrant_api_url(path),
                headers=qdrant_headers(),
                json=payload,
            )
            response.raise_for_status()
            return response.json()
    except Exception as error:
        raise HTTPException(status_code=500, detail=f"Qdrant request failed: {error}") from error


def now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


@contextmanager
def db_conn() -> sqlite3.Connection:
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    try:
        yield conn
        conn.commit()
    finally:
        conn.close()


def verify_token(credentials: HTTPAuthorizationCredentials = Security(security)) -> str:
    if credentials.credentials != API_TOKEN:
        raise HTTPException(status_code=401, detail="Invalid authentication token")
    return credentials.credentials


def init_storage() -> None:
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    FILES_DIR.mkdir(parents=True, exist_ok=True)
    DERIVED_DIR.mkdir(parents=True, exist_ok=True)

    with db_conn() as conn:
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS files (
              file_id TEXT PRIMARY KEY,
              sha256 TEXT NOT NULL,
              filename TEXT NOT NULL,
              mime_type TEXT,
              source_type TEXT NOT NULL,
              area TEXT NOT NULL,
              access_level TEXT NOT NULL,
              owner_type TEXT NOT NULL,
              owner_id TEXT NOT NULL,
              tags TEXT NOT NULL DEFAULT '[]',
              storage_path TEXT NOT NULL,
              parse_status TEXT NOT NULL,
              chunk_count INTEGER NOT NULL DEFAULT 0,
              channel_file_capable INTEGER NOT NULL DEFAULT 1,
              created_at TEXT NOT NULL,
              updated_at TEXT NOT NULL
            )
            """
        )
        conn.execute(
            """
            CREATE UNIQUE INDEX IF NOT EXISTS idx_files_sha_owner_area
            ON files (sha256, owner_id, area)
            """
        )


def ensure_collection() -> None:
    collections_payload = qdrant_request("GET", "/collections")
    collections = collections_payload.get("result", {}).get("collections", [])
    if not any(collection.get("name") == QDRANT_COLLECTION for collection in collections):
        qdrant_request(
            "PUT",
            f"/collections/{QDRANT_COLLECTION}",
            {
                "vectors": {
                    "size": EMBEDDING_DIM,
                    "distance": "Cosine",
                }
            },
        )


def clean_text(text: str) -> str:
    return re.sub(r"\s+", " ", text or "").strip()


def chunk_words(text: str, split_length: int = 500, split_overlap: int = 50) -> list[str]:
    words = text.split()
    if not words:
        return []

    chunks: list[str] = []
    step = max(1, split_length - split_overlap)
    start = 0
    while start < len(words):
        end = min(len(words), start + split_length)
        chunk = " ".join(words[start:end]).strip()
        if chunk:
            chunks.append(chunk)
        if end >= len(words):
            break
        start += step
    return chunks


def embed_text(text: str) -> list[float]:
    if EMBEDDING_PROVIDER == "zhipu":
        api_key = require_env("ZHIPU_API_KEY", ZHIPU_API_KEY)
        url = "https://open.bigmodel.cn/api/paas/v4/embeddings"
        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
        }
        data = {"model": EMBEDDING_MODEL, "input": text}
    elif EMBEDDING_PROVIDER == "openai":
        payload = openai_json_request(
            "/embeddings",
            {"model": EMBEDDING_MODEL, "input": text},
            timeout=30.0,
        )
        return payload["data"][0]["embedding"]
    else:
        raise HTTPException(
            status_code=500,
            detail=f"Unsupported embedding provider: {EMBEDDING_PROVIDER}",
        )

    try:
        with httpx.Client(timeout=30.0) as client:
            response = client.post(url, json=data, headers=headers)
            response.raise_for_status()
            payload = response.json()
            return payload["data"][0]["embedding"]
    except Exception as error:
        raise HTTPException(status_code=500, detail=f"Embedding failed: {error}") from error


def openai_auth_header(explicit_auth: Optional[str] = None) -> str:
    if explicit_auth:
        return explicit_auth
    api_key = require_env("OPENAI_API_KEY", OPENAI_API_KEY)
    return f"Bearer {api_key}"


def openai_json_request(
    path: str,
    payload: dict[str, Any],
    *,
    auth_header: Optional[str] = None,
    timeout: float = 90.0,
) -> dict[str, Any]:
    url = f"{OPENAI_BASE_URL}{path}"
    headers = {
        "Authorization": openai_auth_header(auth_header),
        "Content-Type": "application/json",
    }
    last_error: Optional[Exception] = None

    for attempt in range(1, OPENAI_MAX_RETRIES + 1):
        try:
            response = requests.post(url, json=payload, headers=headers, timeout=timeout)
            response.raise_for_status()
            return response.json()
        except Exception as error:
            last_error = error
            if attempt >= OPENAI_MAX_RETRIES:
                break
            time.sleep((OPENAI_RETRY_BACKOFF_MS / 1000.0) * attempt)

    raise HTTPException(status_code=500, detail=f"OpenAI request failed: {last_error}") from last_error


def understand_image_with_openai(prompt: str, mime_type: str, image_base64: str) -> str:
    payload = {
        "model": OPENAI_VISION_MODEL,
        "messages": [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:{mime_type};base64,{image_base64}"},
                    },
                ],
            }
        ],
        "max_tokens": 1200,
    }
    response = openai_json_request("/chat/completions", payload, timeout=120.0)
    return response["choices"][0]["message"]["content"]


def run_local_vision_ocr(image_path: Path) -> Optional[dict[str, Any]]:
    if platform.system() != "Darwin":
        return None
    if not VISION_OCR_SCRIPT.exists():
        return None

    try:
        completed = subprocess.run(
            ["swift", str(VISION_OCR_SCRIPT), str(image_path)],
            capture_output=True,
            text=True,
            timeout=45,
            check=False,
        )
        if completed.returncode != 0:
            return None
        return json.loads(completed.stdout)
    except Exception:
        return None


def build_ocr_hint(ocr_result: Optional[dict[str, Any]]) -> str:
    if not ocr_result:
        return ""

    full_text = clean_text(str(ocr_result.get("fullText") or ""))
    lines = ocr_result.get("lines") or []
    top_lines = [
        clean_text(str(line.get("text") or ""))
        for line in lines[:10]
        if clean_text(str(line.get("text") or ""))
    ]
    if not full_text and not top_lines:
        return ""

    parts = ["\n\n## OCR 参考\n以下文字由本地 OCR 提取，做聊天截图判断时优先以这些文字为准。"]
    if full_text:
        parts.append(f"\n完整 OCR：\n{full_text}")
    if top_lines:
        parts.append("\n逐行 OCR：\n" + "\n".join(f"- {line}" for line in top_lines))
    return "".join(parts)


def understand_image_bytes(
    image_bytes: bytes,
    filename: str,
    *,
    ocr_result: Optional[dict[str, Any]] = None,
) -> str:
    mime_type = guess_mime_type(filename, None)
    image_base64 = base64.b64encode(image_bytes).decode("utf-8")
    prompt = (
        "请分析这份图片材料，并严格按以下 markdown 结构输出，不要加额外解释。\n\n"
        "## 摘要\n"
        "用 2-4 句话概括材料内容、主题和用途。\n\n"
        "## 可见文字\n"
        "尽量完整提取图片里能看清的文字，没有就写“无”。\n\n"
        "## 角色与方向\n"
        "如果这是聊天截图，请明确说明：\n"
        "- 会话对象是谁\n"
        "- 左侧气泡是谁发的，原文是什么\n"
        "- 右侧气泡是谁发的，原文是什么\n"
        "- 你的判断依据是什么\n"
        "如果不是聊天截图，就写“非聊天截图”。\n\n"
        "## 版面元素\n"
        "列出主要人物、图表、表格、物体、页面结构、场景元素。\n\n"
        "## 关键词\n"
        "- 至少 5 个关键词\n"
    )
    prompt += build_ocr_hint(ocr_result)

    providers = [VISION_PROVIDER]
    if VISION_PROVIDER != "openai" and OPENAI_API_KEY:
        providers.append("openai")

    last_error: Optional[HTTPException] = None

    for provider in providers:
        try:
            if provider == "openai":
                return understand_image_with_openai(prompt, mime_type, image_base64)

            if provider == "zhipu":
                api_key = require_env("ZHIPU_API_KEY", ZHIPU_API_KEY)
                url = "https://open.bigmodel.cn/api/paas/v4/chat/completions"
                headers = {
                    "Authorization": f"Bearer {api_key}",
                    "Content-Type": "application/json",
                }
            elif provider == "gemini":
                api_key = require_env("GEMINI_API_KEY", GEMINI_API_KEY)
                url = f"{GEMINI_BASE_URL}/chat/completions"
                headers = {
                    "Authorization": f"Bearer {api_key}",
                    "Content-Type": "application/json",
                }
            else:
                raise HTTPException(
                    status_code=500,
                    detail=f"Unsupported vision provider: {provider}",
                )

            data = {
                "model": VISION_MODEL,
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {
                                "type": "image_url",
                                "image_url": {"url": f"data:{mime_type};base64,{image_base64}"},
                            },
                        ],
                    }
                ],
            }

            with httpx.Client(timeout=90.0) as client:
                response = client.post(url, json=data, headers=headers)
                response.raise_for_status()
                payload = response.json()
                return payload["choices"][0]["message"]["content"]
        except HTTPException as error:
            last_error = error
        except Exception as error:
            last_error = HTTPException(
                status_code=500,
                detail=f"Image understanding failed: {error}",
            )

    raise last_error or HTTPException(status_code=500, detail="Image understanding failed")


async def proxy_openai_request(request: Request, path: str) -> Response:
    try:
        payload = await request.json()
    except Exception as error:
        raise HTTPException(status_code=400, detail=f"Invalid JSON payload: {error}") from error

    result = await asyncio.to_thread(
        openai_json_request,
        path,
        payload,
        auth_header=request.headers.get("authorization"),
        timeout=180.0,
    )
    return Response(
        content=json.dumps(result, ensure_ascii=False),
        status_code=200,
        headers={"Content-Type": "application/json"},
    )


def guess_mime_type(filename: str, upload_content_type: Optional[str]) -> str:
    guessed, _ = mimetypes.guess_type(filename)
    return upload_content_type or guessed or "application/octet-stream"


def source_type_from_suffix(filename: str) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix in IMAGE_SUFFIXES:
        return "image"
    if suffix in PDF_SUFFIXES:
        return "pdf"
    if suffix in PPT_SUFFIXES:
        return "pptx"
    if suffix in LEGACY_PPT_SUFFIXES:
        return "ppt"
    return "file"


def save_original_file(file_bytes: bytes, sha256_value: str, filename: str) -> Path:
    target_dir = FILES_DIR / sha256_value
    target_dir.mkdir(parents=True, exist_ok=True)
    safe_name = Path(filename).name
    target_path = target_dir / safe_name
    target_path.write_bytes(file_bytes)
    return target_path


def write_derived_chunks(file_id: str, chunks: list[dict[str, Any]]) -> None:
    target_dir = DERIVED_DIR / file_id
    target_dir.mkdir(parents=True, exist_ok=True)
    (target_dir / "chunks.json").write_text(
        json.dumps(chunks, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )


def get_existing_file(sha256_value: str, owner_id: str, area: str) -> Optional[sqlite3.Row]:
    with db_conn() as conn:
        row = conn.execute(
            """
            SELECT * FROM files
            WHERE sha256 = ? AND owner_id = ? AND area = ?
            LIMIT 1
            """,
            (sha256_value, owner_id, area),
        ).fetchone()
    return row


def upsert_file_record(record: dict[str, Any]) -> None:
    with db_conn() as conn:
        conn.execute(
            """
            INSERT INTO files (
              file_id, sha256, filename, mime_type, source_type, area, access_level,
              owner_type, owner_id, tags, storage_path, parse_status, chunk_count,
              channel_file_capable, created_at, updated_at
            )
            VALUES (
              :file_id, :sha256, :filename, :mime_type, :source_type, :area, :access_level,
              :owner_type, :owner_id, :tags, :storage_path, :parse_status, :chunk_count,
              :channel_file_capable, :created_at, :updated_at
            )
            ON CONFLICT(file_id) DO UPDATE SET
              sha256 = excluded.sha256,
              filename = excluded.filename,
              mime_type = excluded.mime_type,
              source_type = excluded.source_type,
              area = excluded.area,
              access_level = excluded.access_level,
              owner_type = excluded.owner_type,
              owner_id = excluded.owner_id,
              tags = excluded.tags,
              storage_path = excluded.storage_path,
              parse_status = excluded.parse_status,
              chunk_count = excluded.chunk_count,
              channel_file_capable = excluded.channel_file_capable,
              updated_at = excluded.updated_at
            """,
            record,
        )


def get_file_record(file_id: str) -> Optional[sqlite3.Row]:
    with db_conn() as conn:
        row = conn.execute(
            "SELECT * FROM files WHERE file_id = ? LIMIT 1",
            (file_id,),
        ).fetchone()
    return row


def row_to_dict(row: sqlite3.Row) -> dict[str, Any]:
    item = dict(row)
    try:
        item["tags"] = json.loads(item.get("tags") or "[]")
    except Exception:
        item["tags"] = []
    item["channel_file_capable"] = bool(item.get("channel_file_capable"))
    return item


def delete_points_for_file(file_id: str) -> None:
    qdrant_request(
        "POST",
        f"/collections/{QDRANT_COLLECTION}/points/delete?wait=true",
        {
            "filter": {
                "must": [
                    {"key": "file_id", "match": {"value": file_id}},
                ]
            }
        },
    )


def upsert_chunks(
    file_id: str,
    filename: str,
    storage_path: str,
    source_type: str,
    area: str,
    access_level: str,
    owner_type: str,
    owner_id: str,
    chunks: list[dict[str, Any]],
) -> None:
    points: list[dict[str, Any]] = []
    for index, chunk in enumerate(chunks):
        text = clean_text(chunk["text"])
        if not text:
            continue
        vector = embed_text(text)
        point = {
            "id": str(uuid.uuid5(uuid.UUID(file_id), f"chunk:{index}")),
            "vector": vector,
            "payload": {
                "file_id": file_id,
                "filename": filename,
                "source_type": source_type,
                "area": area,
                "access_level": access_level,
                "owner_type": owner_type,
                "owner_id": owner_id,
                "page_no": chunk.get("page_no"),
                "slide_no": chunk.get("slide_no"),
                "chunk_no": chunk.get("chunk_no", index),
                "text": text,
                "storage_path": storage_path,
            },
        }
        points.append(point)

    if not points:
        raise HTTPException(status_code=400, detail="No extractable content found")

    qdrant_request(
        "PUT",
        f"/collections/{QDRANT_COLLECTION}/points?wait=true",
        {"points": points},
    )


def process_image(image_path: Path) -> list[dict[str, Any]]:
    ocr_result = run_local_vision_ocr(image_path)
    markdown = understand_image_bytes(
        image_path.read_bytes(),
        image_path.name,
        ocr_result=ocr_result,
    )
    return [{"chunk_no": 0, "text": markdown}]


def render_pdf_page_png(pdf_path: Path, page_index: int) -> bytes:
    if fitz is None:
        raise RuntimeError("PyMuPDF is not installed")
    with fitz.open(pdf_path) as document:
        page = document.load_page(page_index)
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        return pix.tobytes("png")


def process_pdf(pdf_path: Path) -> list[dict[str, Any]]:
    reader = PdfReader(str(pdf_path))
    chunks: list[dict[str, Any]] = []

    for page_index, page in enumerate(reader.pages, start=1):
        extracted = clean_text(page.extract_text() or "")
        if len(extracted) < 40:
            if fitz is not None:
                rendered = render_pdf_page_png(pdf_path, page_index - 1)
                extracted = clean_text(
                    understand_image_bytes(rendered, f"{pdf_path.name}-page-{page_index}.png")
                )

        for chunk_no, chunk in enumerate(chunk_words(extracted), start=1):
            chunks.append(
                {
                    "page_no": page_index,
                    "chunk_no": len(chunks),
                    "text": chunk,
                }
            )

    return chunks


def extract_slide_text(slide: Any) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        text = getattr(shape, "text", "")
        if text:
            cleaned = clean_text(text)
            if cleaned:
                parts.append(cleaned)

    notes_slide = getattr(slide, "notes_slide", None)
    if notes_slide:
        for shape in notes_slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                cleaned = clean_text(text)
                if cleaned:
                    parts.append(f"备注: {cleaned}")

    return "\n".join(parts)


def process_pptx(pptx_path: Path) -> list[dict[str, Any]]:
    presentation = Presentation(str(pptx_path))
    chunks: list[dict[str, Any]] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_text = extract_slide_text(slide)
        if not slide_text:
            continue
        for chunk in chunk_words(slide_text):
            chunks.append(
                {
                    "slide_no": slide_index,
                    "chunk_no": len(chunks),
                    "text": chunk,
                }
            )

    return chunks


def convert_ppt_to_pptx(ppt_path: Path) -> Path:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise HTTPException(
            status_code=500,
            detail="Legacy .ppt conversion requires LibreOffice/soffice, but it is not installed",
        )

    output_dir = Path(tempfile.mkdtemp(prefix="ppt-convert-"))
    try:
        completed = subprocess.run(
            [
                soffice,
                "--headless",
                "--convert-to",
                "pptx",
                "--outdir",
                str(output_dir),
                str(ppt_path),
            ],
            capture_output=True,
            text=True,
            check=False,
        )
        converted = output_dir / f"{ppt_path.stem}.pptx"
        if completed.returncode != 0 or not converted.exists():
            raise HTTPException(
                status_code=500,
                detail=f"Failed to convert .ppt to .pptx: {completed.stderr or completed.stdout or 'unknown error'}",
            )
        return converted
    finally:
        # The converted file is read immediately by the caller, so cleanup can happen there.
        pass


def list_file_hits(file_id: str, limit: int = 5) -> list[dict[str, Any]]:
    scroll_payload = qdrant_request(
        "POST",
        f"/collections/{QDRANT_COLLECTION}/points/scroll",
        {
            "filter": {
                "must": [
                    {"key": "file_id", "match": {"value": file_id}},
                ]
            },
            "limit": limit,
            "with_payload": True,
            "with_vectors": False,
        },
    )
    points = scroll_payload.get("result", {}).get("points", [])

    hits: list[dict[str, Any]] = []
    for point in points:
        payload = point.get("payload") or {}
        hits.append(
            {
                "file_id": payload.get("file_id"),
                "filename": payload.get("filename"),
                "source_type": payload.get("source_type"),
                "page_no": payload.get("page_no"),
                "slide_no": payload.get("slide_no"),
                "chunk_no": payload.get("chunk_no"),
                "excerpt": payload.get("text"),
            }
        )
    return hits


def serialize_search_result(result: Any) -> dict[str, Any]:
    payload = result.get("payload") or {}
    text = payload.get("text") or ""
    return {
        "file_id": payload.get("file_id"),
        "filename": payload.get("filename"),
        "text": text,
        "excerpt": text[:260],
        "score": result.get("score"),
        "area": payload.get("area"),
        "source_type": payload.get("source_type"),
        "owner_type": payload.get("owner_type"),
        "owner_id": payload.get("owner_id"),
        "page_no": payload.get("page_no"),
        "slide_no": payload.get("slide_no"),
        "chunk_no": payload.get("chunk_no"),
        "storage_path": payload.get("storage_path"),
    }


def parse_file_to_chunks(source_type: str, stored_path: Path) -> list[dict[str, Any]]:
    if source_type == "image":
        return process_image(stored_path)
    if source_type == "pdf":
        return process_pdf(stored_path)
    if source_type == "pptx":
        return process_pptx(stored_path)
    if source_type == "ppt":
        converted = convert_ppt_to_pptx(stored_path)
        try:
            return process_pptx(converted)
        finally:
            shutil.rmtree(converted.parent, ignore_errors=True)
    raise HTTPException(status_code=400, detail="Unsupported file type")


init_storage()
ensure_collection()


@app.post("/ingest/file")
async def ingest_file(
    file: UploadFile = File(...),
    area: str = Form(...),
    access_level: str = Form("global"),
    owner_type: str = Form("agent"),
    owner_id: str = Form("default"),
    tags: list[str] = Form(default=[]),
    force_reindex: bool = Form(False),
    _token: str = Depends(verify_token),
) -> dict[str, Any]:
    file_bytes = await file.read()
    filename = file.filename or f"upload-{uuid.uuid4()}"
    suffix = Path(filename).suffix.lower()

    if suffix not in IMAGE_SUFFIXES | PDF_SUFFIXES | PPT_SUFFIXES | LEGACY_PPT_SUFFIXES:
        raise HTTPException(status_code=400, detail="Unsupported file type")

    sha256_value = hashlib.sha256(file_bytes).hexdigest()
    existing = get_existing_file(sha256_value, owner_id, area)
    if existing and not force_reindex:
        return {
            "status": "success",
            "duplicate": True,
            "file_id": existing["file_id"],
            "filename": existing["filename"],
            "type": existing["source_type"],
            "chunks": existing["chunk_count"],
        }

    file_id = existing["file_id"] if existing and force_reindex else str(uuid.uuid4())
    stored_path = save_original_file(file_bytes, sha256_value, filename)
    source_type = source_type_from_suffix(filename)
    mime_type = guess_mime_type(filename, file.content_type)
    chunks = parse_file_to_chunks(source_type, stored_path)

    if existing and force_reindex:
        delete_points_for_file(file_id)

    upsert_chunks(
        file_id=file_id,
        filename=filename,
        storage_path=str(stored_path),
        source_type=source_type,
        area=area,
        access_level=access_level,
        owner_type=owner_type,
        owner_id=owner_id,
        chunks=chunks,
    )
    write_derived_chunks(file_id, chunks)

    timestamp = now_iso()
    created_at = existing["created_at"] if existing else timestamp
    upsert_file_record(
        {
            "file_id": file_id,
            "sha256": sha256_value,
            "filename": filename,
            "mime_type": mime_type,
            "source_type": source_type,
            "area": area,
            "access_level": access_level,
            "owner_type": owner_type,
            "owner_id": owner_id,
            "tags": json.dumps(tags, ensure_ascii=False),
            "storage_path": str(stored_path),
            "parse_status": "ready",
            "chunk_count": len(chunks),
            "channel_file_capable": 1,
            "created_at": created_at,
            "updated_at": timestamp,
        }
    )

    return {
        "status": "success",
        "duplicate": False,
        "file_id": file_id,
        "filename": filename,
        "type": source_type,
        "chunks": len(chunks),
    }


@app.post("/search")
async def search_knowledge(
    query: str = Form(...),
    area: Optional[str] = Form(None),
    access_level: str = Form("global"),
    owner_type: Optional[str] = Form(None),
    owner_id: Optional[str] = Form(None),
    limit: int = Form(5),
    min_score: float = Form(0.35),
    _token: str = Depends(verify_token),
) -> dict[str, Any]:
    query_vector = embed_text(query)

    must_conditions: list[dict[str, Any]] = [
        {"key": "access_level", "match": {"value": access_level}}
    ]
    if area:
        must_conditions.append({"key": "area", "match": {"value": area}})
    if owner_type:
        must_conditions.append({"key": "owner_type", "match": {"value": owner_type}})
    if owner_id:
        must_conditions.append({"key": "owner_id", "match": {"value": owner_id}})

    search_payload = qdrant_request(
        "POST",
        f"/collections/{QDRANT_COLLECTION}/points/search",
        {
            "vector": query_vector,
            "limit": limit,
            "filter": {"must": must_conditions},
            "with_payload": True,
        },
    )
    results = search_payload.get("result", [])
    filtered = [
        serialize_search_result(result)
        for result in results
        if float(result.get("score", 0.0)) >= min_score
    ]

    return {
        "status": "success",
        "query": query,
        "results": filtered,
    }


@app.get("/files/{file_id}")
async def get_file(file_id: str, _token: str = Depends(verify_token)) -> FileResponse:
    row = get_file_record(file_id)
    if not row:
        raise HTTPException(status_code=404, detail="File not found")

    record = row_to_dict(row)
    storage_path = Path(record["storage_path"])
    if not storage_path.exists():
        raise HTTPException(status_code=404, detail="Stored file is missing")

    return FileResponse(
        path=storage_path,
        filename=record["filename"],
        media_type=record.get("mime_type") or "application/octet-stream",
    )


@app.get("/files/{file_id}/meta")
async def get_file_meta(file_id: str, _token: str = Depends(verify_token)) -> dict[str, Any]:
    row = get_file_record(file_id)
    if not row:
        raise HTTPException(status_code=404, detail="File not found")

    return {
        "status": "success",
        "file": row_to_dict(row),
        "hits": list_file_hits(file_id),
    }


@app.post("/openai/v1/embeddings")
async def proxy_openai_embeddings(request: Request) -> Response:
    return await proxy_openai_request(request, "/embeddings")


@app.post("/openai/v1/chat/completions")
async def proxy_openai_chat_completions(request: Request) -> Response:
    return await proxy_openai_request(request, "/chat/completions")


@app.get("/health")
async def health() -> dict[str, Any]:
    return {
        "status": "ok",
        "embedding_provider": EMBEDDING_PROVIDER,
        "model": EMBEDDING_MODEL,
        "vision_provider": VISION_PROVIDER,
        "vision_model": VISION_MODEL,
        "openai_vision_model": OPENAI_VISION_MODEL,
        "collection": QDRANT_COLLECTION,
        "db_path": str(DB_PATH),
    }


# ============================================================================
# Bill Generator Endpoints
# ============================================================================

import sys
sys.path.insert(0, str(RUNTIME_ROOT / "extensions" / "bill-generator-plugin"))

from src.orchestrator import BillGeneratorOrchestrator
from src.ocr import IDCardInfo

# Initialize orchestrator (lazy loading)
_bill_orchestrator: Optional[BillGeneratorOrchestrator] = None


def get_bill_orchestrator() -> BillGeneratorOrchestrator:
    """Get or create bill generator orchestrator."""
    global _bill_orchestrator
    if _bill_orchestrator is None:
        _bill_orchestrator = BillGeneratorOrchestrator()
    return _bill_orchestrator


@app.post("/bill/ocr-id-card")
async def ocr_id_card(
    file: UploadFile = File(...),
    _token: str = Depends(verify_token),
) -> dict[str, Any]:
    """
    OCR recognition for ID card.

    Args:
        file: ID card image file

    Returns:
        Extracted ID card information
    """
    try:
        # Save uploaded file temporarily
        temp_path = Path(tempfile.mktemp(suffix=".jpg"))
        with open(temp_path, "wb") as f:
            content = await file.read()
            f.write(content)

        # Perform OCR
        orchestrator = get_bill_orchestrator()
        id_info: IDCardInfo = orchestrator.recognize_id_card(str(temp_path))

        # Clean up
        temp_path.unlink(missing_ok=True)

        return {
            "name": id_info.name,
            "id_number": id_info.id_number,
            "address": id_info.address,
            "birth_date": id_info.birth_date,
            "confidence": id_info.confidence
        }

    except FileNotFoundError as e:
        raise HTTPException(status_code=404, detail=str(e))
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"OCR failed: {str(e)}")


@app.post("/bill/generate")
async def generate_bill(
    request: Request,
    _token: str = Depends(verify_token),
) -> dict[str, Any]:
    """
    Generate bill PDF from template and identity information.

    Request body:
        {
            "template_path": "/path/to/template.pdf",
            "identity": {
                "name": "张三",
                "id_number": "110101199001011234",
                "address": "北京市东城区..."
            },
            "custom_fields": {
                "金额": 127.35,
                "账期": "2026-03"
            }
        }

    Returns:
        Generation result with output paths and changes
    """
    try:
        body = await request.json()

        template_path = body.get("template_path")
        identity = body.get("identity", {})
        custom_fields = body.get("custom_fields")

        if not template_path:
            raise HTTPException(status_code=400, detail="template_path is required")

        if not identity.get("name"):
            raise HTTPException(status_code=400, detail="identity.name is required")

        # Generate bill
        orchestrator = get_bill_orchestrator()
        result = orchestrator.generate_bill(
            template_path=template_path,
            identity=identity,
            custom_fields=custom_fields
        )

        if not result.success:
            raise HTTPException(status_code=500, detail=result.error)

        return result.to_dict()

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Generation failed: {str(e)}")
